import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Load data
orders = pd.read_excel("northwind_orders.xlsx")
details = pd.read_excel("northwind_orderdetails.xlsx")
customers = pd.read_excel("northwind_customers.xlsx")

# Calculate revenue by country
details["Revenue"] = details["UnitPrice"] * details["Quantity"]
merged = details.merge(orders, on="OrderID").merge(customers, on="CustomerID")
revenue = (
    merged.groupby("Country")["Revenue"]
    .sum()
    .sort_values(ascending=False)
    .reset_index()
)

# Build PowerPoint
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # blank layout
slide = prs.slides.add_slide(slide_layout)

# Title text box
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = txBox.text_frame
tf.text = "Revenue by Country"
tf.paragraphs[0].runs[0].font.size = Pt(28)
tf.paragraphs[0].runs[0].font.bold = True

# Chart data
chart_data = ChartData()
chart_data.categories = list(revenue["Country"])
chart_data.add_series("Revenue", list(revenue["Revenue"]))

# Add bar chart
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5),
    Inches(1.0),
    Inches(9),
    Inches(5.5),
    chart_data,
).chart

chart.has_legend = False
chart.has_title = False

prs.save("revenue_by_country.pptx")
print("Saved revenue_by_country.pptx")
