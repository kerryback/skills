"""Read (and write) a slide deck's speaker notes from its source file.

Two source formats, chosen because each already has a good notes editor built
in: a Quarto reveal.js `.qmd`, where notes live in `::: {.notes}` blocks, and a
PowerPoint `.pptx`, where they live in the notes pane under each slide. The app
does not offer a notes editor of its own — you edit the deck, the app reads it.

Slide *pixels* come from the matching PDF, not from here (see deck.py). What
this module owns is the mapping "slide N has these notes", which is exactly the
thing a PDF cannot tell us.
"""
import re
from pathlib import Path

QMD_EXTS = {".qmd", ".md"}
PPTX_EXTS = {".pptx"}


def detect(path: Path) -> str:
    ext = Path(path).suffix.lower()
    if ext in QMD_EXTS:
        return "qmd"
    if ext in PPTX_EXTS:
        return "pptx"
    if ext == ".ppt":
        raise ValueError(
            "Old-style .ppt isn't supported — open it in PowerPoint and save as "
            ".pptx.")
    if ext == ".pdf":
        raise ValueError(
            "A PDF has no speaker notes. Point the app at the deck you wrote — a "
            ".qmd or .pptx — and it will use the PDF alongside it for the slide "
            "images.")
    raise ValueError(f"Unsupported deck source '{ext}'. Use a .qmd or a .pptx.")


def read_slides(path: Path) -> list:
    """Per-slide {index, title, slide_text, notes} in presentation order."""
    path = Path(path)
    kind = detect(path)
    slides = (parse_qmd(path.read_text(encoding="utf-8")) if kind == "qmd"
              else parse_pptx(path))
    for i, s in enumerate(slides):
        s["index"] = i
    return slides


# --------------------------------------------------------------------------- #
# Quarto / reveal.js
# --------------------------------------------------------------------------- #
_FENCE = re.compile(r"^(`{3,}|~{3,})(.*)$")
_DIV = re.compile(r"^(:{3,})\s*(.*)$")
# A heading needs a space after the hashes, or nothing at all: `## Title` and a
# bare `##` (an untitled slide, which decks use for a full-bleed or continuation
# slide) are both headings; `#hashtag` is not.
_HEADING = re.compile(r"^(#{1,6})(?:\s+(.*))?$")
_HRULE = re.compile(r"^(?:-{3,}|\*{3,}|_{3,})\s*$")
_ATTRS = re.compile(r"\{[^}]*\}\s*$")


def _front_matter(lines: list) -> tuple:
    """(parsed front matter, index of the first body line)."""
    if not lines or lines[0].strip() != "---":
        return {}, 0
    for i in range(1, len(lines)):
        if lines[i].strip() in ("---", "..."):
            block = "\n".join(lines[1:i])
            try:
                import yaml
                meta = yaml.safe_load(block)
            except Exception:
                meta = None
            return (meta if isinstance(meta, dict) else {}), i + 1
    return {}, 0


def _slide_level(meta: dict) -> int:
    """Heading depth that starts a new slide. Quarto's default is 2 (`##`), and
    `#` then makes a section-title slide — which is still a slide, and still a
    PDF page, so anything at or above this level breaks a slide."""
    fmt = meta.get("format")
    if isinstance(fmt, dict):
        for key, val in fmt.items():
            if "revealjs" in str(key) and isinstance(val, dict):
                if val.get("slide-level") is not None:
                    return int(val["slide-level"])
    if meta.get("slide-level") is not None:
        return int(meta["slide-level"])
    return 2


def _clean_title(text: str) -> str:
    text = _ATTRS.sub("", text).strip()
    text = re.sub(r"[*_`]", "", text)
    return text.strip()


def _slide_text(lines: list) -> str:
    """The slide's visible content, flattened. Only used to give the drafting
    agent context and to label thumbnails, so light de-markdowning is enough."""
    out = []
    for ln in lines:
        s = ln.strip()
        if not s:
            continue
        s = re.sub(r"^[-*+]\s+|^\d+[.)]\s+|^>\s?|^#+\s+", "", s)
        s = _ATTRS.sub("", s)
        s = re.sub(r"[*_`]", "", s)
        if s.strip():
            out.append(s.strip())
    return re.sub(r"\s+", " ", " ".join(out)).strip()


def parse_qmd(text: str) -> list:
    """Split a Quarto reveal deck into slides, pulling out each one's notes.

    Slide breaks follow Quarto: a heading at or above `slide-level`, a `---`
    rule, and — when the front matter sets a title — the title slide Quarto
    generates before any of them. Content before the first heading belongs to
    that title slide, which is where a title slide's notes have to go.
    """
    lines = text.splitlines()
    meta, start = _front_matter(lines)
    level = _slide_level(meta)

    slides = []

    def new_slide(title=""):
        slides.append({"title": title, "body": [], "notes": [],
                       "notes_span": None})
        return slides[-1]

    if meta.get("title"):
        new_slide(_clean_title(str(meta["title"])))

    cur = slides[-1] if slides else None
    fence = None          # open code fence marker, or None
    div_depth = 0
    notes_at = None       # div depth at which the current .notes block opened
    prev_blank = True

    for lineno in range(start, len(lines)):
        line = lines[lineno].rstrip()
        stripped = line.strip()

        # --- fenced code: everything inside is literal ---
        fm = _FENCE.match(stripped)
        if fm:
            tok = fm.group(1)
            if fence is None:
                fence = tok
            elif tok[0] == fence[0] and len(tok) >= len(fence) and not fm.group(2).strip():
                fence = None
            if cur and notes_at is None:
                cur["body"].append(line)
            prev_blank = False
            continue
        if fence is not None:
            if cur and notes_at is None:
                cur["body"].append(line)
            prev_blank = False
            continue

        # --- pandoc divs: `::: {.notes}` is the one we care about ---
        dm = _DIV.match(stripped)
        if dm:
            attrs = dm.group(2).strip()
            if attrs:
                if notes_at is None and re.search(r"(^|[.{\s])\.?notes\b", attrs):
                    if cur is None:
                        cur = new_slide()
                    notes_at = div_depth
                    cur["notes_span"] = [lineno, lineno]
                    div_depth += 1
                    prev_blank = False
                    continue
                div_depth += 1
            else:
                div_depth = max(0, div_depth - 1)
                if notes_at is not None and div_depth == notes_at:
                    cur["notes_span"][1] = lineno
                    notes_at = None
                    prev_blank = False
                    continue
            if cur and notes_at is None:
                cur["body"].append(line)
            prev_blank = False
            continue

        if notes_at is not None:
            cur["notes"].append(line)
            continue

        hm = _HEADING.match(stripped)
        if hm and len(hm.group(1)) <= level:
            cur = new_slide(_clean_title(hm.group(2) or ""))
            prev_blank = False
            continue

        if _HRULE.match(stripped) and prev_blank:
            cur = new_slide()
            prev_blank = False
            continue

        if cur is None:
            # Blank lines between the front matter and the first heading are not
            # a slide. Only real content before the first heading is (Quarto
            # puts it on the title slide, or makes an untitled one).
            if not stripped:
                prev_blank = True
                continue
            cur = new_slide()
        cur["body"].append(line)
        prev_blank = not stripped

    return [{
        "index": i,
        "title": s["title"],
        "slide_text": _slide_text(s["body"]),
        "notes": "\n".join(s["notes"]).strip(),
        "notes_span": s["notes_span"],
    } for i, s in enumerate(slides)]


def write_qmd_notes(path: Path, notes_by_index: dict) -> int:
    """Set the notes on the given slides of a .qmd, in place.

    Replaces an existing `::: {.notes}` block or appends one to the slide. Only
    the slides named in `notes_by_index` are touched; everything else in the
    file — including formatting the parser doesn't model — is left byte for
    byte alone. Returns the number of slides changed.
    """
    path = Path(path)
    lines = path.read_text(encoding="utf-8").splitlines()
    slides = parse_qmd("\n".join(lines))

    # Slide boundaries, so an appended block lands inside the right slide. A
    # slide ends where the next one begins; the last runs to end of file.
    starts = _slide_starts(lines, slides)

    # Apply from the bottom up so earlier edits don't shift later line numbers.
    changed = 0
    for i in sorted(notes_by_index, reverse=True):
        if i >= len(slides):
            continue
        text = (notes_by_index[i] or "").strip()
        span = slides[i]["notes_span"]
        block = ["::: {.notes}", *text.splitlines(), ":::"] if text else []
        if span:
            lines[span[0]:span[1] + 1] = block
        else:
            end = starts[i + 1] if i + 1 < len(starts) else len(lines)
            while end > starts[i] and not lines[end - 1].strip():
                end -= 1
            lines[end:end] = ([""] + block) if block else []
        changed += 1

    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return changed


def _slide_starts(lines: list, slides: list) -> list:
    """First line index of each slide, recovered by re-walking the headings and
    rules the parser used. The front-matter title slide starts at the body."""
    meta, start = _front_matter(lines)
    level = _slide_level(meta)
    starts = []
    if meta.get("title"):
        starts.append(start)
    fence = None
    div_depth = 0
    notes_at = None
    prev_blank = True
    for i in range(start, len(lines)):
        stripped = lines[i].strip()
        fm = _FENCE.match(stripped)
        if fm:
            tok = fm.group(1)
            if fence is None:
                fence = tok
            elif tok[0] == fence[0] and len(tok) >= len(fence) and not fm.group(2).strip():
                fence = None
            prev_blank = False
            continue
        if fence is not None:
            prev_blank = False
            continue
        dm = _DIV.match(stripped)
        if dm:
            attrs = dm.group(2).strip()
            if attrs:
                if notes_at is None and re.search(r"(^|[.{\s])\.?notes\b", attrs):
                    notes_at = div_depth
                div_depth += 1
            else:
                div_depth = max(0, div_depth - 1)
                if notes_at is not None and div_depth == notes_at:
                    notes_at = None
            prev_blank = False
            continue
        if notes_at is not None:
            continue
        hm = _HEADING.match(stripped)
        if hm and len(hm.group(1)) <= level:
            starts.append(i)
            prev_blank = False
            continue
        if _HRULE.match(stripped) and prev_blank:
            starts.append(i)
            prev_blank = False
            continue
        if not starts and stripped:
            starts.append(i)
        prev_blank = not stripped
    return starts[:len(slides)] or [start]


# --------------------------------------------------------------------------- #
# PowerPoint
# --------------------------------------------------------------------------- #
def _pptx_presentation(path: Path):
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover - dependency is installed by the launcher
        raise RuntimeError(
            "Reading .pptx needs the python-pptx package (pip install python-pptx).")
    return Presentation(str(path))


def parse_pptx(path: Path) -> list:
    prs = _pptx_presentation(path)
    out = []
    for i, slide in enumerate(prs.slides):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
        except (AttributeError, ValueError):
            title = ""
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = (shape.text_frame.text or "").strip()
                if t:
                    texts.append(t)
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        out.append({
            "index": i,
            "title": title.splitlines()[0] if title else "",
            "slide_text": re.sub(r"\s+", " ", " ".join(texts)).strip(),
            "notes": notes,
            "hidden": slide._element.get("show") == "0",
        })
    return out


def count_hidden_pptx(path: Path) -> int:
    """Hidden slides, which PowerPoint leaves out of an exported PDF by default —
    the most likely reason a .pptx and its PDF disagree on how many slides there
    are."""
    try:
        return sum(1 for s in parse_pptx(path) if s.get("hidden"))
    except Exception:
        return 0


def write_pptx_notes(path: Path, notes_by_index: dict) -> int:
    """Set the notes on the given slides of a .pptx, in place. Creates a notes
    slide for any slide that has none. Returns the number of slides changed."""
    prs = _pptx_presentation(path)
    slides = list(prs.slides)
    changed = 0
    for i, text in notes_by_index.items():
        if i < 0 or i >= len(slides):
            continue
        slides[i].notes_slide.notes_text_frame.text = (text or "").strip()
        changed += 1
    if changed:
        prs.save(str(path))
    return changed


def write_notes(path: Path, notes_by_index: dict) -> int:
    """Write notes back into whichever source format this is."""
    path = Path(path)
    kind = detect(path)
    notes_by_index = {int(k): v for k, v in notes_by_index.items()}
    if kind == "qmd":
        return write_qmd_notes(path, notes_by_index)
    return write_pptx_notes(path, notes_by_index)
